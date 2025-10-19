from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
NAVY = RGBColor(28, 42, 74)
GREEN = RGBColor(39, 174, 96)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 242, 245)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Navy background
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
    text_frame = subtitle_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = GREEN
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White background
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    
    # Title bar
    title_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)

# Slide 1: Title
add_title_slide(prs, "VAT-Aware Invoice Template Pack", "SARS-Compliant Invoicing for South African SMMEs")

# Slide 2: What's Included
add_content_slide(prs, "What's Included in This Pack", [
    "📄 Full Tax Invoice Template (for transactions > R5,000)",
    "📄 Abridged Tax Invoice Template (for transactions ≤ R5,000)",
    "📄 Credit Note Template (for returns and adjustments)",
    "📄 Quotation Template (professional quote format)",
    "📚 Comprehensive VAT Invoicing Guide (30+ pages)",
    "✅ All templates are SARS-compliant and ready to use"
])

# Slide 3: Why This Matters
add_content_slide(prs, "Why VAT-Compliant Invoicing Matters", [
    "⚖️ Legal Requirement: VAT-registered businesses must issue compliant invoices",
    "💰 Get Paid Faster: Professional invoices build trust and speed payment",
    "🛡️ Avoid Penalties: Non-compliance can result in fines up to 200% of tax",
    "📊 Claim Input VAT: Proper invoices let you recover VAT on purchases",
    "🏢 B2B Advantage: Corporate clients require VAT-compliant documentation"
])

# Slide 4: Understanding VAT in SA
add_content_slide(prs, "VAT Registration in South Africa", [
    "Current VAT Rate: 15% (increased from 14% in April 2018)",
    "Mandatory Registration: Turnover > R1 million in 12 months",
    "Voluntary Registration: Turnover > R50,000 in 12 months",
    "Benefits: Claim input VAT, professional credibility, B2B advantage",
    "Obligations: File returns monthly or bi-monthly, keep records for 5 years"
])

# Slide 5: Full vs Abridged
add_content_slide(prs, "Full vs. Abridged Tax Invoices", [
    "Full Tax Invoice (> R5,000):",
    "  • Must include customer's VAT number (if registered)",
    "  • Detailed breakdown of VAT amount",
    "  • All 10 mandatory SARS elements required",
    "",
    "Abridged Tax Invoice (≤ R5,000):",
    "  • Simplified format for quick invoicing",
    "  • Customer's VAT number not required",
    "  • Perfect for retail and cash sales"
])

# Slide 6: How to Use the Templates
add_content_slide(prs, "How to Use the Templates", [
    "1️⃣ Customize with your business details (name, address, VAT number)",
    "2️⃣ Set up sequential invoice numbering (INV-001, INV-002, etc.)",
    "3️⃣ Add customer information and line items",
    "4️⃣ Ensure VAT is calculated at 15%",
    "5️⃣ Include your banking details for payment",
    "6️⃣ Save each invoice with a unique filename",
    "7️⃣ Keep copies for 5 years (SARS requirement)"
])

# Slide 7: SA-Specific Tips
add_content_slide(prs, "South African Best Practices", [
    "⚡ Load Shedding: Use cloud storage for invoices (accessible during outages)",
    "⏰ Late Payments: SA average is 60-90 days—offer early payment discounts",
    "🔒 Banking Security: Warn customers about banking detail scams",
    "💳 Payment Gateways: Use PayFast, Yoco, or Ozow for instant payments",
    "📱 Mobile Invoicing: Abridged invoices work great on smartphones",
    "🎯 Follow-Up: Send reminders at 7, 14, and 21 days"
])

# Slide 8: Common Mistakes to Avoid
add_content_slide(prs, "Common Mistakes to Avoid", [
    "❌ Using wrong invoice type (Abridged for > R5,000)",
    "❌ Incorrect VAT calculation (always 15% of VAT-exclusive amount)",
    "❌ Missing customer's VAT number (required for Full Tax Invoices)",
    "❌ Duplicate invoice numbers (each must be unique and sequential)",
    "❌ Not issuing credit notes for overcharges",
    "❌ Charging VAT before registration (illegal!)"
])

# Slide 9: Next Steps
add_content_slide(prs, "Your Next Steps", [
    "✅ Customize the templates with your business details",
    "✅ Read the comprehensive guide (02-VAT-Invoice-Guide.pdf)",
    "✅ Set up your invoice numbering system",
    "✅ Create a filing system (digital and/or physical)",
    "✅ Consider invoicing software (Xero, QuickBooks, Zoho)",
    "✅ Consult a tax practitioner if you have questions",
    "",
    "📞 SARS Contact Centre: 0800 00 7277"
])

# Save
prs.save('/home/ubuntu/Operiva/artifacts/01-README-VAT-Invoice.pptx')
print("✓ Created comprehensive README presentation (9 slides)")


#!/usr/bin/env python3.11
"""
Insert the SWOT to TOWS process flowchart into SWOT Analysis Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load the presentation
prs = Presentation('/home/ubuntu/Operiva/artifacts/04-SWOT-Analysis-Presentation.pptx')

# Get the blank layout
blank_layout = None
for layout in prs.slide_layouts:
    if 'blank' in layout.name.lower() or layout.name == 'Blank':
        blank_layout = layout
        break

if blank_layout is None:
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

# Insert new slide after slide 4 (index 4, so it becomes slide 5)
slide = prs.slides.add_slide(blank_layout)

# Move the new slide to position 4 (0-indexed, so it becomes slide 5 in 1-indexed)
xml_slides = prs.slides._sldIdLst
slides_list = list(xml_slides)
xml_slides.remove(slides_list[-1])
xml_slides.insert(4, slides_list[-1])

# Add title
left = Inches(0.5)
top = Inches(0.3)
width = Inches(9)
height = Inches(0.7)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = "From SWOT Analysis to Actionable Strategies"

# Format title
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(30)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(28, 42, 74)  # Navy blue
title_paragraph.alignment = PP_ALIGN.CENTER

# Add flowchart image
img_path = '/home/ubuntu/Operiva/flowcharts/swot-to-tows-process.png'
left = Inches(0.5)
top = Inches(1.2)
width = Inches(9)
slide.shapes.add_picture(img_path, left, top, width=width)

# Save the updated presentation
prs.save('/home/ubuntu/Operiva/artifacts/04-SWOT-Analysis-Presentation-UPDATED.pptx')

print("✓ Flowchart inserted into SWOT Analysis Presentation (new slide 5)")
print(f"  Total slides: {len(prs.slides)}")
print("  Saved as: 04-SWOT-Analysis-Presentation-UPDATED.pptx")


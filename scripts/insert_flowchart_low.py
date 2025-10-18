#!/usr/bin/env python3.11
"""
Insert the Levels of Work 4-Levels flowchart into LOW Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load the presentation
prs = Presentation('/home/ubuntu/Operiva/artifacts/04-Levels-of-Work-Presentation.pptx')

# Get the blank layout
blank_layout = None
for layout in prs.slide_layouts:
    if 'blank' in layout.name.lower() or layout.name == 'Blank':
        blank_layout = layout
        break

if blank_layout is None:
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

# Insert new slide after slide 3 (index 3, so it becomes slide 4)
slide = prs.slides.add_slide(blank_layout)

# Move the new slide to position 3 (0-indexed, so it becomes slide 4 in 1-indexed)
xml_slides = prs.slides._sldIdLst
slides_list = list(xml_slides)
xml_slides.remove(slides_list[-1])
xml_slides.insert(3, slides_list[-1])

# Add title
left = Inches(0.5)
top = Inches(0.3)
width = Inches(9)
height = Inches(0.7)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = "The 4-Level Organizational Hierarchy"

# Format title
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(28)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(28, 42, 74)  # Navy blue
title_paragraph.alignment = PP_ALIGN.CENTER

# Add flowchart image
img_path = '/home/ubuntu/Operiva/flowcharts/levels-of-work-4-levels.png'
left = Inches(0.3)
top = Inches(1.1)
width = Inches(9.4)
slide.shapes.add_picture(img_path, left, top, width=width)

# Save the updated presentation
prs.save('/home/ubuntu/Operiva/artifacts/04-Levels-of-Work-Presentation-UPDATED.pptx')

print("✓ Flowchart inserted into Levels of Work Presentation (new slide 4)")
print(f"  Total slides: {len(prs.slides)}")
print("  Saved as: 04-Levels-of-Work-Presentation-UPDATED.pptx")


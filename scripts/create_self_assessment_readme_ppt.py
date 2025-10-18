#!/usr/bin/env python3.11
"""
Create README presentation for Owner-Manager Self-Assessment
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define Operiva colors
NAVY = RGBColor(28, 42, 74)
GREEN = RGBColor(39, 174, 96)
LIGHT_GRAY = RGBColor(245, 245, 245)
DARK_GRAY = RGBColor(100, 100, 100)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.shapes.add_shape(
        1,  # Rectangle
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = NAVY
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.2), Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = GREEN
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    
    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.8), Inches(8.4), Inches(5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for item in content_items:
        p = content_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)
        p.level = 0
    
    return slide

# Slide 1: Title
add_title_slide(
    prs,
    "Owner-Manager Self-Assessment",
    "The Capability Compass for SA SMMEs"
)

# Slide 2: What's Inside
add_content_slide(
    prs,
    "📦 What's in This Package?",
    [
        "✅ This README (quick overview)",
        "✅ Comprehensive Guide (25+ page PDF)",
        "✅ Interactive Assessment Worksheet (Excel)",
        "✅ Presentation Deck (visual summary)",
        "",
        "⏱️ Time Required: 20-30 minutes",
        "🎯 Purpose: Identify your capability strengths and gaps",
        "📊 Output: Personalized development roadmap"
    ]
)

# Slide 3: Why This Matters
add_content_slide(
    prs,
    "❓ Why Self-Assessment?",
    [
        "70-80% of SA SMMEs fail within 5 years",
        "Primary cause: Capability gaps in owner-managers",
        "",
        "Traditional assessment tools are:",
        "  • Text-heavy and boring",
        "  • Designed for corporates, not SMMEs",
        "  • Culturally misaligned",
        "",
        "The Capability Compass uses scenario-based assessment:",
        "  ✅ Real SA SMME situations",
        "  ✅ Practical, not theoretical",
        "  ✅ Immediate, actionable feedback"
    ]
)

# Slide 4: The 5 Capabilities
add_content_slide(
    prs,
    "🧭 The 5 Capability Areas",
    [
        "1️⃣ Financial Management",
        "   Cash flow, pricing, bookkeeping, VAT compliance",
        "",
        "2️⃣ Strategic Planning & Vision",
        "   Goal-setting, SWOT, market analysis, adaptability",
        "",
        "3️⃣ Leadership & People",
        "   Delegation, team management, culture, succession",
        "",
        "4️⃣ Marketing & Sales",
        "   Customer acquisition, branding, digital presence",
        "",
        "5️⃣ Operations & Technology",
        "   Processes, efficiency, automation, POPIA compliance"
    ]
)

# Slide 5: How to Use
add_content_slide(
    prs,
    "📝 How to Use This Tool",
    [
        "STEP 1: Open 03-Self-Assessment-Worksheet.xlsx",
        "",
        "STEP 2: Read each scenario carefully",
        "  • 10 scenarios per capability area (50 total)",
        "  • Choose the response that matches what YOU would do",
        "  • Be honest - this is for YOUR benefit!",
        "",
        "STEP 3: Review your Dashboard",
        "  • See scores for each capability (out of 30)",
        "  • Overall score (out of 150)",
        "  • Color-coded ratings and priority actions",
        "",
        "STEP 4: Create an action plan",
        "  • Focus on lowest-scoring areas",
        "  • Download relevant Operiva artifacts",
        "  • Set 3 goals for next 90 days"
    ]
)

# Slide 6: Scoring Guide
add_content_slide(
    prs,
    "📊 Understanding Your Scores",
    [
        "Each Capability Area (out of 30 points):",
        "  🌟 27-30 = STRONG (best practice level)",
        "  ✅ 21-26 = COMPETENT (solid foundation)",
        "  ⚠️ 12-20 = DEVELOPING (significant gaps)",
        "  🚨 0-11 = URGENT (critical weaknesses)",
        "",
        "Overall Score (out of 150 points):",
        "  120-150 = Strong owner-manager capabilities",
        "  90-119 = Competent, prioritize 1-2 areas",
        "  60-89 = Developing, seek targeted support",
        "  0-59 = Urgent, consider coaching/mentorship"
    ]
)

# Slide 7: What Makes This Different
add_content_slide(
    prs,
    "✨ What Makes This Different?",
    [
        "🇿🇦 SA-Specific Scenarios:",
        "  • Load shedding impact",
        "  • VAT and SARS compliance",
        "  • Township vs suburban markets",
        "  • LSM targeting, stokvels, WhatsApp marketing",
        "",
        "📱 Mobile-Friendly:",
        "  • Use Excel app on your phone",
        "  • Minimal text, visual feedback",
        "",
        "🎯 Actionable:",
        "  • Not just a score - get specific next steps",
        "  • Links to relevant Operiva artifacts",
        "  • 90-day development plan template"
    ]
)

# Slide 8: Next Steps
add_content_slide(
    prs,
    "🚀 Your Next Steps",
    [
        "1. Read the Guide (02-Self-Assessment-Guide.pdf)",
        "   Understand the 5 capabilities and why they matter",
        "",
        "2. Complete the Worksheet (03-Self-Assessment-Worksheet.xlsx)",
        "   20-30 minutes, be honest with yourself",
        "",
        "3. Review Your Dashboard",
        "   Identify your top 2 development areas",
        "",
        "4. Take Action",
        "   Download Operiva artifacts for your weak areas",
        "   Set 3 specific goals for next 90 days",
        "",
        "5. Retake in 6 Months",
        "   Track your progress and celebrate growth!"
    ]
)

# Slide 9: Resources
add_content_slide(
    prs,
    "📚 Recommended Operiva Artifacts",
    [
        "If Financial Management is weak:",
        "  → Cash Flow Forecast, VAT Compliance Checklist",
        "",
        "If Strategic Planning is weak:",
        "  → SWOT Analysis Pack, Market Entry Playbook",
        "",
        "If Leadership is weak:",
        "  → Levels of Work Organogram, Delegation Guide",
        "",
        "If Marketing is weak:",
        "  → Digital Marketing Playbook, Customer Journey Map",
        "",
        "If Operations is weak:",
        "  → Process Documentation Template, POPIA Compliance"
    ]
)

# Save presentation
output_path = '/home/ubuntu/Operiva/artifacts/01-README-Self-Assessment.pptx'
prs.save(output_path)

print(f"✅ README presentation created: {output_path}")
print(f"   - 9 slides")
print(f"   - Operiva branding (navy + green)")


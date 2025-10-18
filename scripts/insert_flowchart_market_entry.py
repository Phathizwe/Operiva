#!/usr/bin/env python3.11
"""
Insert the 4-phase process flowchart into Market Entry Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load the presentation
prs = Presentation('/home/ubuntu/Operiva/artifacts/04-Market-Entry-Presentation.pptx')

# Get the blank layout (usually index 6, but let's find it)
blank_layout = None
for layout in prs.slide_layouts:
    if 'blank' in layout.name.lower() or layout.name == 'Blank':
        blank_layout = layout
        break

# If no blank layout found, use the first one
if blank_layout is None:
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

# Insert new slide after slide 2 (index 2, so it becomes slide 3)
slide = prs.slides.add_slide(blank_layout)

# Move the new slide to position 2 (0-indexed, so it becomes slide 3 in 1-indexed)
xml_slides = prs.slides._sldIdLst
slides_list = list(xml_slides)
# Move last slide to position 2
xml_slides.remove(slides_list[-1])
xml_slides.insert(2, slides_list[-1])

# Add title
left = Inches(0.5)
top = Inches(0.4)
width = Inches(9)
height = Inches(0.8)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = "The 4-Phase Market Entry Process"

# Format title
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(32)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(28, 42, 74)  # Navy blue #1C2A4A
title_paragraph.alignment = PP_ALIGN.CENTER

# Add flowchart image
img_path = '/home/ubuntu/Operiva/flowcharts/market-entry-4-phase-process.png'
left = Inches(0.3)
top = Inches(1.4)
width = Inches(9.4)
slide.shapes.add_picture(img_path, left, top, width=width)

# Save the updated presentation
prs.save('/home/ubuntu/Operiva/artifacts/04-Market-Entry-Presentation-UPDATED.pptx')

print("✓ Flowchart inserted into Market Entry Presentation (new slide 3)")
print(f"  Total slides: {len(prs.slides)}")
print("  Saved as: 04-Market-Entry-Presentation-UPDATED.pptx")

